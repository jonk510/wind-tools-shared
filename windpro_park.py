#!/usr/bin/env python3
"""
WindPRO PARK Results Summary — park_summary.py

Core logic: extract data from WindPRO PARK PDF exports and build a
PowerPoint presentation.  Works as a standalone CLI script and as an
importable module (used by the Streamlit app in app.py).

Slide structure:
  • Cover slide (uses Cover_Renewables_Wind template layout)
  • One slide per calculation — satellite map + key metrics + per-WTG wake loss chart
  • Final slide — comparison table across all calculations

CLI usage:
    python park_summary.py                                  # auto-discover PARK_*.pdf
    python park_summary.py a.pdf b.pdf                      # specific files
    python park_summary.py --template template.pptx a.pdf   # branded template
"""

import io
import re
import sys
from pathlib import Path

import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.colors as mcolors
import matplotlib.patches as mpatches
import matplotlib.patheffects as mpe
import matplotlib.pyplot as plt
from matplotlib.path import Path as MplPath
import pdfplumber

_NOISE_IMPORT_ERR = None
try:
    from shared.acoustics import (
        compute_noise_grid as _compute_noise_grid,
        _build_elev_interp,
    )
    from shared.srtm import fetch_srtm_elevation as _fetch_srtm_elevation
    HAS_NOISE = True
except Exception as _e:
    HAS_NOISE = False
    _NOISE_IMPORT_ERR = str(_e)

try:
    import fitz  # pymupdf
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    import contextily as cx
    from pyproj import Transformer
    HAS_CTX = True
except ImportError:
    HAS_CTX = False

try:
    import geopandas as gpd
    HAS_GPD = True
except ImportError:
    HAS_GPD = False

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from lxml import etree

# ── Default loss assumptions ──────────────────────────────────────────────────
DEFAULT_LOSSES: dict = {
    'Temp derating loss [%]':           None,
    'Availability loss [%]':            3.5,
    'Electrical loss [%]':              2.2,
    'Turbine performance loss [%]':     2.0,
    'Degradation [%]':                  1.5,
}


def load_shapefile(path: str):
    """
    Load a shapefile or zipped shapefile (.zip containing .shp/.dbf/.shx).
    Returns a GeoDataFrame, or None if geopandas is unavailable or loading fails.
    """
    if not HAS_GPD:
        return None
    try:
        return gpd.read_file(path)
    except Exception:
        return None


def _ordinal(n: int) -> str:
    """Return ordinal string for integer n, e.g. 1 → '1st', 3 → '3rd'."""
    suffix = ('th' if 11 <= n % 100 <= 13
              else {1: 'st', 2: 'nd', 3: 'rd'}.get(n % 10, 'th'))
    return f'{n}{suffix}'


# ── Colour palette — extracted directly from template01.pptx ─────────────────
DARK_NAVY  = RGBColor(0x0F, 0x24, 0x3E)  # table label-col bg  (#0F243E)
NAVY       = RGBColor(0x24, 0x45, 0x7E)  # accent1             (#24457E)
ORANGE     = RGBColor(0xEB, 0x9E, 0x05)  # gold highlight      (#EB9E05)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GREY       = RGBColor(0x80, 0x80, 0x80)  # secondary / loss data
BORDER_CLR = RGBColor(0xA2, 0xA2, 0xA2)  # cell border
FONT_NAME  = 'Arial'

# ── Template layout geometry (10" × 5.62") ────────────────────────────────────
HDR_H   = 0.94   # master header image height [inches]
FOOT_Y  = 5.26   # conservative footer boundary [inches]
LM      = 0.41   # left margin
RM      = 9.59   # right edge


# ─────────────────────────────────────────────────────────────────────────────
# Data extraction
# ─────────────────────────────────────────────────────────────────────────────

def _last_number(text: str) -> float | None:
    nums = re.findall(r'[\d,]+\.?\d*', text)
    return float(nums[-1].replace(',', '')) if nums else None


def _min_spacing_d(wtg_coords: dict, rotor_m: float) -> float | None:
    """Minimum inter-WTG distance as a multiple of rotor diameter."""
    if len(wtg_coords) < 2 or not rotor_m:
        return None
    import math
    pts = list(wtg_coords.values())
    mean_lat = math.radians(sum(p[1] for p in pts) / len(pts))
    m_per_lon = 111_320 * math.cos(mean_lat)
    m_per_lat = 111_320
    min_sq = float('inf')
    for i in range(len(pts)):
        for j in range(i + 1, len(pts)):
            dx = (pts[j][0] - pts[i][0]) * m_per_lon
            dy = (pts[j][1] - pts[i][1]) * m_per_lat
            d = dx * dx + dy * dy
            if d < min_sq:
                min_sq = d
    return math.sqrt(min_sq) / rotor_m


def extract(pdf_path: str) -> dict:
    """
    Extract PARK metrics from a WindPRO PDF export.

    Handles two main result table formats:
      Format A (with "other losses" column):
        … NAME  net_MWh  wake_free_MWh  wake_%  free_ws  reduced_ws
      Format B (without "other losses" column):
        … NAME  result_MWh  wake_%  free_ws  reduced_ws

    Returns a dict with summary scalars, 'wtgs' list, and 'wtg_coords' dict.
    """
    d: dict = {'wtgs': [], 'wtg_coords': {}}

    with pdfplumber.open(pdf_path) as pdf:
        pages = {i + 1: (p.extract_text() or '') for i, p in enumerate(pdf.pages)}
        full  = '\n'.join(pages.values())

        # ── Calculation name & date ─────────────────────────────────────────
        m = re.search(r'Calculation:\s*(.+?)(?:\n|WTG:)', full)
        d['calc_name'] = m.group(1).strip() if m else Path(pdf_path).stem

        m = re.search(r'Calculated:\s*\n(\d{1,2}/\d{2}/\d{4})', full)
        d['calc_date'] = m.group(1) if m else ''

        # ── Total capacity + turbine count (Time varying AEP section) ───────
        m = re.search(r'([\d.]+) MW based on (\d+) turbines', full)
        if m:
            d['total_mw'] = float(m.group(1))
            d['num_wtgs'] = int(m.group(2))

        # ── Detect per-WTG table format ─────────────────────────────────────
        main_result_text = pages.get(1, '') + pages.get(2, '') + pages.get(3, '')
        has_other_losses = 'other losses' in main_result_text.lower()

        # ── Parse per-WTG rows ───────────────────────────────────────────────
        # Take LAST 4 (Format B) or 5 (Format A) numbers from the full line to
        # avoid picking up numbers embedded in model names or Name fields.
        wtg_row_re = re.compile(r'^\s*(\d+)(?:Yes|No)\s+')

        manuf = model = ''
        rated_kw = rotor_m = hub_m = 0.0

        for page_num in range(1, min(len(pages) + 1, 8)):
            page_text = pages.get(page_num, '')
            if 'Reference WTGs' in page_text and 'Main Result' not in page_text:
                break
            for line in page_text.splitlines():
                wm = wtg_row_re.match(line)
                if not wm:
                    continue
                wtg_num  = int(wm.group(1))
                all_nums = [float(n.replace(',', ''))
                            for n in re.findall(r'[\d,]+\.?\d*', line)]

                if has_other_losses and len(all_nums) >= 5:
                    mwh_a, mwh_b, wake_pct, free_ws, reduced_ws = all_nums[-5:]
                    wake_free_mwh = max(mwh_a, mwh_b)
                    net_mwh       = min(mwh_a, mwh_b)
                elif not has_other_losses and len(all_nums) >= 4:
                    result_mwh, wake_pct, free_ws, reduced_ws = all_nums[-4:]
                    wake_free_mwh = net_mwh = result_mwh
                else:
                    continue

                if not (0 <= wake_pct <= 50 and 2 <= free_ws <= 30 and 2 <= reduced_ws <= 30):
                    continue

                if not manuf:
                    kw_m = re.search(r'(?<= )(\d{1,2},\d{3})(?= )', line)
                    if kw_m:
                        rated_kw  = float(kw_m.group(1).replace(',', ''))
                        diam_nums = re.findall(r'\d+\.\d+', line[kw_m.end():])
                        if len(diam_nums) >= 2:
                            rotor_m = float(diam_nums[0])
                            hub_m   = float(diam_nums[1])
                        before_kw = line[wm.end():kw_m.start()].strip()
                        parts = before_kw.split()
                        model = parts[-1] if parts else ''
                        manuf = ' '.join(parts[:-1]) if len(parts) > 1 else (parts[0] if parts else '')

                d['wtgs'].append({
                    'wtg_num':       wtg_num,
                    'wake_free_mwh': wake_free_mwh,
                    'net_mwh':       net_mwh,
                    'wake_loss_pct': wake_pct,
                    'free_ws':       free_ws,
                    'reduced_ws':    reduced_ws,
                })

        d.update({
            'manufacturer': manuf,
            'wtg_model':    model,
            'rated_mw':     round(rated_kw / 1000, 1) if rated_kw else None,
            'rotor_m':      rotor_m,
            'hub_m':        hub_m,
        })

        if d['wtgs']:
            d['mean_free_ws']    = round(
                sum(w['free_ws']    for w in d['wtgs']) / len(d['wtgs']), 1)
            d['mean_reduced_ws'] = round(
                sum(w['reduced_ws'] for w in d['wtgs']) / len(d['wtgs']), 1)

        # ── Production Analysis (page varies; scan dynamically) ──────────────
        prod = ''
        for page_num, text in pages.items():
            if ('Production Analysis' in text
                    and 'new WTGs' in text
                    and 'Model based energy' in text):
                prod = text
                break

        m = re.search(r'Model based energy \[MWh\](.*)', prod)
        if m:
            d['gross_aep_mwh'] = _last_number(m.group(1))

        m = re.search(r'Resulting energy \[MWh\](.*)', prod)
        if m:
            d['park_yield_mwh'] = _last_number(m.group(1))

        m = re.search(r'Decrease due to wake losses \[%\]([\d.\s]+)', prod)
        if m:
            nums = re.findall(r'\d+\.\d+', m.group(1))
            d['wake_loss_pct'] = float(nums[-1]) if nums else None

        # ── WTG coordinates (lon/lat) from coordinate/description pages ──────
        coord_re = re.compile(
            r'(\d+)(?:New|Yes|No)?\s+([\d.]+)°\s*E\s+([-\d.]+)°\s*N'
        )
        # Collect only new-WTG numbers (from the energy table)
        new_wtg_nums = {w['wtg_num'] for w in d['wtgs']}
        for page_text in pages.values():
            if '°' not in page_text:
                continue
            for m in coord_re.finditer(page_text):
                num = int(m.group(1))
                if num in new_wtg_nums and num not in d['wtg_coords']:
                    d['wtg_coords'][num] = (float(m.group(2)), float(m.group(3)))

        # ── UTM fallback: coordinate table with merged Easting/Northing ─────
        # Some exports use a UTM Easting/Southing table.  pdfplumber merges
        # adjacent columns without spaces, so easting and northing appear
        # concatenated: "1New 350,9496,690,869297.2ENVISION…"
        # Easting is always 6 digits (###,###), southing 7 digits (#,###,###).
        if not d['wtg_coords']:
            try:
                from pyproj import Transformer as _Tr
                zone_m = re.search(r'UTM[^\n]*Zone:\s*(\d+)', full, re.IGNORECASE)
                utm_zone = int(zone_m.group(1)) if zone_m else 50
                to_wgs84 = _Tr.from_crs(f'EPSG:{32700 + utm_zone}',
                                        'EPSG:4326', always_xy=True)

                # Merged UTM table rows (coordinate table pages)
                utm_row_re = re.compile(
                    r'(\d+)New\s+(\d{3},\d{3})(\d,\d{3},\d{3})'
                )
                # Spaced UTM (per-WTG wind-data pages)
                utm_spaced_re = re.compile(r'East:\s*([\d,]+)\s+North:\s*([\d,]+)')
                wtg_hdr_re    = re.compile(r'Wind data:\s*(\d+)\s+-')

                for page_text in pages.values():
                    # Merged table format
                    if 'Easting' in page_text or 'Southing' in page_text:
                        for m in utm_row_re.finditer(page_text):
                            num = int(m.group(1))
                            if num not in new_wtg_nums or num in d['wtg_coords']:
                                continue
                            e = float(m.group(2).replace(',', ''))
                            n = float(m.group(3).replace(',', ''))
                            lon, lat = to_wgs84.transform(e, n)
                            d['wtg_coords'][num] = (lon, lat)
                    # Per-WTG page format
                    wm = wtg_hdr_re.search(page_text)
                    um = utm_spaced_re.search(page_text)
                    if wm and um:
                        num = int(wm.group(1))
                        if num in new_wtg_nums and num not in d['wtg_coords']:
                            e = float(um.group(1).replace(',', ''))
                            n = float(um.group(2).replace(',', ''))
                            lon, lat = to_wgs84.transform(e, n)
                            d['wtg_coords'][num] = (lon, lat)
            except Exception:
                pass

        # ── Renumber WTGs north → south ───────────────────────────────────
        # Only when all WTGs have coords; lat descending = north first.
        if d['wtg_coords'] and len(d['wtg_coords']) == len(d['wtgs']):
            ordered = sorted(d['wtg_coords'], key=lambda n: d['wtg_coords'][n][1],
                             reverse=True)
            remap = {old: new for new, old in enumerate(ordered, 1)}
            d['wtg_coords'] = {remap[k]: v for k, v in d['wtg_coords'].items()}
            for w in d['wtgs']:
                if w['wtg_num'] in remap:
                    w['wtg_num'] = remap[w['wtg_num']]

        # ── Minimum inter-WTG spacing ──────────────────────────────────────
        if d['wtg_coords'] and d.get('rotor_m'):
            d['min_spacing_d'] = _min_spacing_d(d['wtg_coords'], d['rotor_m'])

    return d


def apply_losses(d: dict, losses: dict) -> dict:
    """Calculate p50_aep_mwh from park_yield_mwh × (1 − each loss)."""
    if 'park_yield_mwh' in d:
        f = 1.0
        for v in losses.values():
            if v:
                f *= 1.0 - v / 100.0
        d['p50_aep_mwh'] = d['park_yield_mwh'] * f
    return d


# ─────────────────────────────────────────────────────────────────────────────
# Visual assets (matplotlib)
# ─────────────────────────────────────────────────────────────────────────────

def _pdf_page_bytes(pdf_path: str, dpi: int = 200) -> bytes | None:
    """Render the last PDF page (WindPRO map) as PNG bytes via pymupdf."""
    if not HAS_FITZ:
        return None
    doc = fitz.open(pdf_path)
    page = doc[-1]
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat)
    return pix.tobytes('png')


def _make_wtg_marker() -> MplPath:
    """3-bladed rotor plan-view marker — tapered blades + hub circle."""
    def _blade(angle_deg):
        a = np.radians(angle_deg)
        c, s = np.cos(a), np.sin(a)
        pts = np.array([
            [ 0.00,  0.00],
            [-0.10,  0.18],
            [-0.05,  1.00],
            [ 0.05,  1.00],
            [ 0.10,  0.18],
            [ 0.00,  0.00],
        ])
        rot = np.array([[c, -s], [s, c]])
        return (pts @ rot.T).tolist()

    verts, codes = [], []
    for angle in [90, 210, 330]:
        blade = _blade(angle)
        verts += blade
        codes += [MplPath.MOVETO] + [MplPath.LINETO] * (len(blade) - 2) + [MplPath.CLOSEPOLY]

    t = np.linspace(0, 2 * np.pi, 16, endpoint=False)
    hub = np.column_stack([0.16 * np.cos(t), 0.16 * np.sin(t)])
    verts += hub.tolist() + [hub[0].tolist()]
    codes += [MplPath.MOVETO] + [MplPath.LINETO] * (len(hub) - 1) + [MplPath.CLOSEPOLY]

    return MplPath(verts, codes)


_WTG_MARKER = _make_wtg_marker()

# ── Noise contour colours (green → yellow → orange → red → purple) ────────────
_NOISE_COLOURS = [
    '#27ae60', '#a9d65d', '#f9ca24', '#f0932b', '#eb4d4b', '#8e44ad', '#2d3436',
]


def _noise_cmap_norm(levels):
    n      = len(levels) - 1
    colors = (_NOISE_COLOURS * ((n // len(_NOISE_COLOURS)) + 1))[:n]
    cmap   = mcolors.ListedColormap(colors)
    norm   = mcolors.BoundaryNorm(levels, len(colors))
    return cmap, norm


def compute_noise_overlay(wtg_coords: dict, hub_height: float, Lw_bands: dict,
                           resolution: int = 150, buffer_m: float = 3000.0,
                           hr: float = 4.0, G: float = 0.5,
                           terrain_xyz=None,
                           use_terrain: bool = False,
                           use_shielding: bool = False) -> dict | None:
    """
    Compute a noise grid from WGS84 WTG coordinates.

    Auto-detects the UTM zone so distances are accurate, optionally loads
    SRTM terrain or a user-supplied XYZ DataFrame, then reprojects the
    finished UTM grid to Web Mercator for satellite-map overlay.

    Parameters
    ----------
    terrain_xyz   : pd.DataFrame with X, Y, Z columns (projected to epsg_utm)
                    or None.  When None and use_terrain=True, SRTM is downloaded.
    use_terrain   : if False, flat terrain is assumed (fast).
    use_shielding : apply ISO 9613-2 §8 terrain barrier attenuation.
    """
    if not HAS_NOISE or not HAS_CTX or not wtg_coords:
        return None
    try:
        from scipy.interpolate import griddata as _griddata

        lons = [v[0] for v in wtg_coords.values()]
        lats = [v[1] for v in wtg_coords.values()]

        # ── Auto-detect UTM zone ──────────────────────────────────────────
        clon = sum(lons) / len(lons)
        clat = sum(lats) / len(lats)
        zone = int((clon + 180) / 6) + 1
        epsg_utm = (32700 if clat < 0 else 32600) + zone

        # ── WGS84 → UTM ───────────────────────────────────────────────────
        to_utm  = Transformer.from_crs('EPSG:4326', f'EPSG:{epsg_utm}', always_xy=True)
        xs_utm, ys_utm = to_utm.transform(lons, lats)
        wtg_xy  = np.column_stack([xs_utm, ys_utm])

        # ── Terrain ───────────────────────────────────────────────────────
        if terrain_xyz is not None:
            get_elev  = _build_elev_interp(terrain_xyz)
            wtg_elevs = get_elev(wtg_xy)
        elif use_terrain:
            xyz_df    = _fetch_srtm_elevation(
                wtg_xy, epsg_utm, buffer_m=buffer_m + 2000, grid_n=40)
            get_elev  = _build_elev_interp(xyz_df)
            wtg_elevs = get_elev(wtg_xy)
        else:
            get_elev  = None
            wtg_elevs = np.zeros(len(wtg_xy))

        # ── UTM noise grid ────────────────────────────────────────────────
        xmin = min(xs_utm) - buffer_m;  xmax = max(xs_utm) + buffer_m
        ymin = min(ys_utm) - buffer_m;  ymax = max(ys_utm) + buffer_m
        xx_utm, yy_utm = np.meshgrid(np.linspace(xmin, xmax, resolution),
                                     np.linspace(ymin, ymax, resolution))
        grid_pts  = np.column_stack([xx_utm.ravel(), yy_utm.ravel()])
        elev_grid = (get_elev(grid_pts).reshape(xx_utm.shape)
                     if get_elev is not None else np.zeros_like(xx_utm))

        noise_utm = _compute_noise_grid(
            wtg_xy, wtg_elevs, Lw_bands, hub_height,
            xx_utm, yy_utm, elev_grid,
            hr=hr, G=G, use_shielding=use_shielding)

        # ── Reproject UTM → Web Mercator for satellite map overlay ────────
        to_merc     = Transformer.from_crs(f'EPSG:{epsg_utm}', 'EPSG:3857', always_xy=True)
        xm, ym      = to_merc.transform(xx_utm.ravel(), yy_utm.ravel())
        xx_merc, yy_merc = np.meshgrid(np.linspace(xm.min(), xm.max(), resolution),
                                        np.linspace(ym.min(), ym.max(), resolution))
        noise_merc  = _griddata(
            np.column_stack([xm, ym]),
            noise_utm.ravel(),
            np.column_stack([xx_merc.ravel(), yy_merc.ravel()]),
            method='linear',
        ).reshape(xx_merc.shape)

        return {'xx': xx_merc, 'yy': yy_merc, 'noise_grid': noise_merc,
                'epsg_utm': epsg_utm}
    except Exception:
        return None


def _render_wtg_map(ax, fig, wtg_coords: dict, wtgs: list[dict],
                    xs: list, ys: list, dpi: int,
                    rotor_m: float = 0, coord_is_meters: bool = True,
                    shapes: list | None = None,
                    noise_overlay: dict | None = None) -> bytes:
    """
    Shared render: scatter WTG markers + labels + 3D circles onto ax/fig,
    then return PNG bytes.  xs/ys are already in the axes coordinate system
    (either Web Mercator or plain lon/lat).
    """
    wake_by_num = {w['wtg_num']: w['wake_loss_pct'] for w in wtgs}
    net_by_num  = {w['wtg_num']: w['net_mwh']       for w in wtgs}
    n_wtgs      = len(wtg_coords)

    wakes   = [wake_by_num.get(n, 0) for n in wtg_coords]
    max_w   = max(wakes) if max(wakes) > 0 else 1
    cmap    = plt.cm.RdYlGn_r
    colours = [cmap(w / max_w) for w in wakes]

    span_x  = max(xs) - min(xs) if len(xs) > 1 else 1
    span_y  = max(ys) - min(ys) if len(ys) > 1 else 1
    area    = span_x * span_y if span_x > 0 else 1
    s       = max(150, min(600, int(40_000_000 / max(area, 1) * n_wtgs)))

    # Shapefile overlays — white 20% fill / 50% border, above basemap, below WTGs
    if HAS_GPD and shapes:
        target_crs = 'EPSG:3857' if coord_is_meters else 'EPSG:4326'
        for gdf in shapes:
            try:
                gdf_proj = gdf.to_crs(target_crs) if gdf.crs is not None else gdf
                gdf_proj.plot(ax=ax,
                              facecolor=(1.0, 1.0, 1.0, 0.20),
                              edgecolor=(1.0, 1.0, 1.0, 0.50),
                              linewidth=1.2,
                              zorder=7)
            except Exception:
                pass

    # Noise contour overlay — zorder 6, above basemap/shapes, below WTG markers
    if noise_overlay and coord_is_meters:
        try:
            lvls           = noise_overlay.get('contour_levels', [35, 40, 45])
            cmap_n, norm_n = _noise_cmap_norm(lvls)
            xx_n  = noise_overlay['xx']
            yy_n  = noise_overlay['yy']
            ng    = noise_overlay['noise_grid']

            # Coloured contour lines only (no fill)
            from matplotlib.lines import Line2D as _Line2D
            colours_n   = [cmap_n(norm_n(lv + 0.01)) for lv in lvls]
            legend_lines = []
            for lv, col in zip(lvls, colours_n):
                try:
                    ax.contour(xx_n, yy_n, ng, levels=[lv],
                               colors=[col], linewidths=1.4, alpha=0.95, zorder=7)
                    legend_lines.append(
                        _Line2D([0], [0], color=col, linewidth=1.4,
                                label=f'{lv:g} dB(A)'))
                except Exception:
                    pass

            # Inline labels
            cl = ax.contour(xx_n, yy_n, ng, levels=lvls,
                            colors='white', linewidths=0.0, alpha=0.0, zorder=7)
            ax.clabel(cl, fmt='%g dB(A)', fontsize=4.5, inline=True,
                      manual=False, use_clabeltext=True)

            if legend_lines:
                ax.legend(handles=legend_lines, loc='lower right',
                          fontsize=4.5, framealpha=0.6,
                          title='Noise', title_fontsize=4.5,
                          facecolor='#1e2a3a', labelcolor='white',
                          edgecolor='white')
        except Exception:
            pass

    ax.scatter(xs, ys,
               marker=_WTG_MARKER, s=s, c=colours,
               edgecolors='white', linewidths=0.8, zorder=10)

    # 3D exclusion circles (radius = 3 × rotor diameter)
    if rotor_m > 0:
        r = 3 * rotor_m if coord_is_meters else 3 * rotor_m / 111_320
        xlim, ylim = ax.get_xlim(), ax.get_ylim()
        for x, y in zip(xs, ys):
            ax.add_patch(mpatches.Circle(
                (x, y), r, fill=False,
                edgecolor='white', linewidth=0.6, linestyle='--',
                alpha=0.5, zorder=9))
        ax.set_xlim(xlim); ax.set_ylim(ylim)  # circles must not expand view
        stroke_leg = [mpe.withStroke(linewidth=1.5, foreground='black')]
        ax.text(0.02, 0.03, f'-- 3D zone  (r = {3 * rotor_m:.0f} m)',
                transform=ax.transAxes, fontsize=4.0,
                color='white', alpha=0.8, path_effects=stroke_leg)

    fsize     = 5.5 if n_wtgs <= 25 else 4.5
    stroke    = [mpe.withStroke(linewidth=2, foreground='black')]
    line_h    = fsize * 1.3            # points per line
    gwh_y     = 5 - line_h             # y-offset for GWh (line 2)
    wake_y    = 5 - 2 * line_h         # y-offset for wake% (line 3)

    for num, (x, y) in zip(wtg_coords.keys(), zip(xs, ys)):
        net  = net_by_num.get(num, 0)
        wake = wake_by_num.get(num, 0)
        # T# in white
        ax.annotate(f"T{num}", (x, y),
                    xytext=(5, 5), textcoords='offset points',
                    fontsize=fsize, color='white', fontweight='bold',
                    path_effects=stroke, zorder=11)
        # GWh in light blue
        ax.annotate(f"{net/1000:.1f} GWh", (x, y),
                    xytext=(5, gwh_y), textcoords='offset points',
                    fontsize=fsize, color='#7EC8E3', fontweight='bold',
                    path_effects=stroke, zorder=11)
        # Wake % in orange
        ax.annotate(f"{wake:.1f}% wake", (x, y),
                    xytext=(5, wake_y), textcoords='offset points',
                    fontsize=fsize, color='#EB9E05', fontweight='bold',
                    path_effects=stroke, zorder=11)

    ax.axis('off')
    fig.patch.set_facecolor('black')
    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches=None, facecolor='black')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def _plain_map_bytes(wtg_coords: dict, wtgs: list[dict],
                     width_in: float, height_in: float, dpi: int,
                     rotor_m: float = 0, shapes: list | None = None,
                     noise_overlay: dict | None = None) -> bytes | None:
    """
    Fallback map — WTG symbols on a dark background using raw lon/lat
    (no pyproj projection).  Used when contextily/pyproj are unavailable
    or satellite tile fetch fails.
    """
    if not wtg_coords:
        return None
    try:
        lons = [v[0] for v in wtg_coords.values()]
        lats = [v[1] for v in wtg_coords.values()]
        pad_lon = max((max(lons) - min(lons)) * 0.15, 0.005)
        pad_lat = max((max(lats) - min(lats)) * 0.15, 0.005)

        fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
        ax.set_facecolor('#1e2a3a')
        ax.set_xlim(min(lons) - pad_lon, max(lons) + pad_lon)
        ax.set_ylim(min(lats) - pad_lat, max(lats) + pad_lat)
        ax.set_aspect('auto')
        return _render_wtg_map(ax, fig, wtg_coords, wtgs, lons, lats, dpi,
                               rotor_m=rotor_m, coord_is_meters=False,
                               shapes=shapes, noise_overlay=None)
    except Exception:
        return None


def satellite_map_bytes(wtg_coords: dict, wtgs: list[dict],
                        width_in: float, height_in: float,
                        dpi: int = 150, rotor_m: float = 0,
                        shapes: list | None = None,
                        noise_overlay: dict | None = None) -> bytes | None:
    """
    Satellite basemap (ESRI World Imagery) with 3-bladed WTG symbols
    coloured green→red by wake loss.
    Falls back to a plain dark-background map if tiles or projection fail.
    Returns None only when no WTG coordinates are available at all.
    """
    if not wtg_coords:
        return None

    if not HAS_CTX:
        return _plain_map_bytes(wtg_coords, wtgs, width_in, height_in, dpi,
                                rotor_m=rotor_m, shapes=shapes, noise_overlay=noise_overlay)

    try:
        to_merc = Transformer.from_crs('EPSG:4326', 'EPSG:3857', always_xy=True)
        lons = [v[0] for v in wtg_coords.values()]
        lats = [v[1] for v in wtg_coords.values()]
        xs_all, ys_all = to_merc.transform(lons, lats)

        span_x = max(xs_all) - min(xs_all) if len(xs_all) > 1 else 1
        span_y = max(ys_all) - min(ys_all) if len(ys_all) > 1 else 1
        pad_x  = max(span_x * 0.12, 2000)
        pad_y  = max(span_y * 0.12, 2000)

        fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=dpi)
        ax.set_xlim(min(xs_all) - pad_x, max(xs_all) + pad_x)
        ax.set_ylim(min(ys_all) - pad_y, max(ys_all) + pad_y)
        ax.set_aspect('auto')

        try:
            cx.add_basemap(ax, crs='EPSG:3857',
                           source=cx.providers.Esri.WorldImagery,
                           attribution=False, zoom='auto')
        except Exception:
            ax.set_facecolor('#1e2a3a')

        xs_wtg, ys_wtg = to_merc.transform(lons, lats)
        return _render_wtg_map(ax, fig, wtg_coords, wtgs,
                               list(xs_wtg), list(ys_wtg), dpi,
                               rotor_m=rotor_m, coord_is_meters=True,
                               shapes=shapes, noise_overlay=noise_overlay)

    except Exception:
        # Projection pipeline failed — fall back to plain lon/lat map
        return _plain_map_bytes(wtg_coords, wtgs, width_in, height_in, dpi,
                                rotor_m=rotor_m, shapes=shapes, noise_overlay=None)


def map_image_bytes(wtg_coords: dict, wtgs: list[dict],
                    width_in: float, height_in: float,
                    rotor_m: float = 0,
                    shapes: list | None = None,
                    noise_overlay: dict | None = None) -> bytes | None:
    """Custom WTG map (satellite or plain dark). No PDF fallback ever."""
    return satellite_map_bytes(wtg_coords, wtgs, width_in, height_in,
                               rotor_m=rotor_m, shapes=shapes,
                               noise_overlay=noise_overlay)


def wake_chart_bytes(wtgs: list[dict], avg_pct: float,
                     width_in: float, height_in: float) -> bytes:
    """Per-WTG wake loss bar chart, green→red gradient. Returns PNG bytes."""
    nums   = [w['wtg_num']       for w in wtgs]
    wakes  = [w['wake_loss_pct'] for w in wtgs]
    max_w  = max(wakes) if wakes else 1
    colours = [plt.cm.RdYlGn_r(v / max_w) for v in wakes]

    fig, ax = plt.subplots(figsize=(width_in, height_in), dpi=150)
    ax.bar(nums, wakes, color=colours, edgecolor='white', linewidth=0.3, zorder=3)
    ax.axhline(avg_pct, color='#35496A', linewidth=1.2, linestyle='--',
               label=f'Avg  {avg_pct:.1f}%', zorder=4)
    ax.legend(fontsize=6.5, framealpha=0.85, loc='upper right')
    ax.set_xlabel('WTG', fontsize=7)
    ax.set_ylabel('Wake Loss [%]', fontsize=7)
    ax.tick_params(labelsize=6.5)
    ax.grid(axis='y', alpha=0.3, linewidth=0.4, zorder=0)
    ax.set_facecolor('#F4F6F9')
    ax.set_xlim(min(nums) - 0.8, max(nums) + 0.8)
    ax.set_ylim(0, max_w * 1.22)
    fig.patch.set_facecolor('white')
    fig.tight_layout(pad=0.4)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# python-pptx helpers
# ─────────────────────────────────────────────────────────────────────────────

def _layout_by_name(prs: Presentation, *names: str):
    """Return the first layout whose name matches any of the given names."""
    for name in names:
        for layout in prs.slide_layouts:
            if layout.name == name:
                return layout
    return prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]


def _set_slide_title(slide, text: str, font_size: int = 11):
    """Write text into the title placeholder (idx=0)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)
            return
    # Fallback: free textbox if no title placeholder
    tb  = slide.shapes.add_textbox(Inches(LM), Inches(0.28), Inches(9.0), Inches(0.35))
    p   = tb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(font_size)
    run.font.bold      = True
    run.font.color.rgb = NAVY


_A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _tc_pr(cell):
    """Return <a:tcPr>, inserting it at position 0 in <a:tc> if absent."""
    tc = cell._tc
    pr = tc.find(f'{{{_A_NS}}}tcPr')
    if pr is None:
        pr = etree.Element(f'{{{_A_NS}}}tcPr')
        tc.insert(0, pr)
    return pr


def _apply_cell_borders(cell, color=BORDER_CLR, width_pt=1.0):
    """Solid 1pt border on all 4 sides of a table cell.

    Must be called before cell.fill.solid() — OOXML schema requires <a:lnL/R/T/B>
    to precede <a:solidFill> in <a:tcPr>, otherwise PowerPoint ignores the borders.
    """
    pr  = _tc_pr(cell)
    w   = str(int(width_pt * 12700))
    col = str(color)          # RGBColor.__str__ → 'RRGGBB'
    for tag in ('lnL', 'lnR', 'lnT', 'lnB'):
        old = pr.find(f'{{{_A_NS}}}{tag}')
        if old is not None:
            pr.remove(old)
        ln = etree.SubElement(pr, f'{{{_A_NS}}}{tag}',
                              w=w, cap='flat', cmpd='sng', algn='ctr')
        sf = etree.SubElement(ln, f'{{{_A_NS}}}solidFill')
        etree.SubElement(sf, f'{{{_A_NS}}}srgbClr', val=col)
        etree.SubElement(ln, f'{{{_A_NS}}}prstDash', val='solid')
        etree.SubElement(ln, f'{{{_A_NS}}}round')
        etree.SubElement(ln, f'{{{_A_NS}}}headEnd', type='none', w='med', len='med')
        etree.SubElement(ln, f'{{{_A_NS}}}tailEnd', type='none', w='med', len='med')


def _set_cell(cell, text: str, bg=None, fg=BLACK,
              bold=False, italic=False, size=8,
              align=PP_ALIGN.CENTER):
    # Margins + borders must be written before fill: OOXML requires <a:lnL/R/T/B>
    # to appear before <a:solidFill> in <a:tcPr>.
    pr = _tc_pr(cell)
    pr.set('marL', '57150')
    pr.set('marR', '57150')
    pr.set('marT', '38100')
    pr.set('marB', '25400')
    pr.set('anchor', 'ctr')

    _apply_cell_borders(cell)

    if bg is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

    tf = cell.text_frame
    tf.word_wrap = True
    cell.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        r = p.runs[0]
        r.font.name      = FONT_NAME
        r.font.size      = Pt(size)
        r.font.bold      = bold
        r.font.italic    = italic
        r.font.color.rgb = fg


def _mini_label(slide, text: str, x, y, w, h=Inches(0.20)):
    """Small bold navy section label."""
    tb  = slide.shapes.add_textbox(x, y, w, h)
    tf  = tb.text_frame
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Pt(0)
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(7.5)
    run.font.bold      = True
    run.font.color.rgb = NAVY


def _hmerge_row(tbl, row_idx: int, col_start: int, col_end: int):
    """Merge table cells horizontally from col_start to col_end (0-indexed)."""
    n_span = col_end - col_start + 1
    if n_span <= 1:
        return
    first_tc = tbl.cell(row_idx, col_start)._tc
    first_tc.set('gridSpan', str(n_span))
    for ci in range(col_start + 1, col_end + 1):
        tc = tbl.cell(row_idx, ci)._tc
        tc.set('hMerge', '1')
        # Clear text content of hidden cells
        txBody = tc.find(f'{{{_A_NS}}}txBody')
        if txBody is not None:
            for p_elem in txBody.findall(f'{{{_A_NS}}}p'):
                for r_elem in list(p_elem.findall(f'{{{_A_NS}}}r')):
                    p_elem.remove(r_elem)


# ─────────────────────────────────────────────────────────────────────────────
# Cover slide
# ─────────────────────────────────────────────────────────────────────────────

def add_cover_slide(prs: Presentation, title: str,
                    subtitle: str = '', subsubtitle: str = ''):
    """Add a branded cover slide — title placeholder + two textbox tiers."""
    layout = _layout_by_name(prs,
                             'Cover_Renewables_Wind',
                             'Cover_Renewables_Solar',
                             'Cover_No Image',
                             'Cover_City Image',
                             'Cover_Regional Image')
    slide = prs.slides.add_slide(layout)

    # ── Tier 1: title placeholder (idx=0) ────────────────────────────────
    title_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True
            title_ph = ph
            break

    # Compute sub-heading positions relative to title placeholder bottom
    H = prs.slide_height.inches
    if title_ph is not None:
        title_bot = (title_ph.top + title_ph.height) / 914400
        sub_x = title_ph.left / 914400
        sub_w = title_ph.width / 914400
    else:
        title_bot = 3.8
        sub_x, sub_w = LM, RM - LM

    sub_y    = min(title_bot + 0.05, H - 0.90) - 0.591  # shift up 15 mm clear of template line
    subsub_y = min(sub_y + 0.44, H - 0.42)

    # ── Tier 2: subtitle (orange, bold) ──────────────────────────────────
    if subtitle:
        tb = slide.shapes.add_textbox(
            Inches(sub_x), Inches(sub_y), Inches(sub_w), Inches(0.44))
        tf = tb.text_frame
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        run = tf.paragraphs[0].add_run()
        run.text           = subtitle
        run.font.name      = FONT_NAME
        run.font.size      = Pt(20)
        run.font.bold      = True
        run.font.color.rgb = ORANGE

    # ── Tier 3: sub-subtitle (orange, smaller) ────────────────────────────
    if subsubtitle:
        tb = slide.shapes.add_textbox(
            Inches(sub_x), Inches(subsub_y), Inches(sub_w), Inches(0.35))
        tf = tb.text_frame
        tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
        run = tf.paragraphs[0].add_run()
        run.text           = subsubtitle
        run.font.name      = FONT_NAME
        run.font.size      = Pt(14)
        run.font.color.rgb = ORANGE


# ─────────────────────────────────────────────────────────────────────────────
# Per-calculation slide
# ─────────────────────────────────────────────────────────────────────────────

def add_calc_slide(prs: Presentation, d: dict, pdf_path: str,
                   blank_layout, losses: dict, shapes: list | None = None,
                   noise_overlay: dict | None = None):
    """
    Proportional layout respecting the template header/footer zones:
      Left  ~48%  – Key Metrics table  +  Per-WTG Wake Loss chart
      Right ~50%  – Satellite map (with WTG markers) or WindPRO PDF map

    Content starts at y = HDR_H + gap to avoid the master header image.
    """
    W = prs.slide_width.inches
    H = prs.slide_height.inches

    # Content zone — safely inside header/footer
    cx0  = LM                     # left edge  = 0.41"
    cy0  = HDR_H + 0.08           # top edge   ≈ 1.02"
    cw   = RM - LM                # total width = 9.18"
    ch   = FOOT_Y - cy0           # total height ≈ 4.24"

    # Panel widths
    left_w  = cw * 0.475          # ≈ 4.36"
    gap     = cw * 0.020          # ≈ 0.18"
    right_x = cx0 + left_w + gap  # ≈ 4.97"
    right_w = RM - right_x        # ≈ 4.62"

    slide = prs.slides.add_slide(blank_layout)
    _set_slide_title(slide, f"PARK Summary  ·  {d.get('calc_name', '')}")

    # ── Right panel: satellite map (never PDF) ────────────────────────────
    img = map_image_bytes(d.get('wtg_coords', {}), d.get('wtgs', []),
                          right_w, ch, rotor_m=d.get('rotor_m') or 0,
                          shapes=shapes, noise_overlay=noise_overlay)
    if img:
        pic = slide.shapes.add_picture(
            io.BytesIO(img),
            Inches(right_x), Inches(cy0),
            width=Inches(right_w), height=Inches(ch))
        # Remove the template's default thick border from the picture shape.
        # Two sources: direct <a:ln> in spPr, and theme lnRef in <p:style>.
        _PML_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        spPr = pic._element.spPr
        for old_ln in list(spPr.findall(f'{{{_A_NS}}}ln')):
            spPr.remove(old_ln)
        ln = etree.SubElement(spPr, f'{{{_A_NS}}}ln', w='0')
        etree.SubElement(ln, f'{{{_A_NS}}}noFill')
        style = pic._element.find(f'{{{_PML_NS}}}style')
        if style is not None:
            lnRef = style.find(f'{{{_A_NS}}}lnRef')
            if lnRef is not None:
                lnRef.set('idx', '0')
    else:
        ph = slide.shapes.add_textbox(
            Inches(right_x), Inches(cy0), Inches(right_w), Inches(ch))
        ph.text_frame.text = 'No WTG coordinates found in PDF\n(satellite map unavailable)'

    # ── Left panel: metrics ───────────────────────────────────────────────
    lbl_h    = 0.20
    tbl_h    = ch * 0.535
    gap_mid  = 0.08
    chart_lbl_y = cy0 + lbl_h + tbl_h + gap_mid
    chart_h  = FOOT_Y - chart_lbl_y - lbl_h - 0.05

    p50 = d.get('p50_aep_mwh', 0)
    mw  = d.get('total_mw', 0)
    n   = d.get('num_wtgs', 0)
    flh = round(p50 / mw) if mw else 0
    cf  = f"{p50 / (mw * 8760) * 100:.1f}%" if mw else '-'

    min_sp = d.get('min_spacing_d')
    metrics = [
        ('Date',            d.get('calc_date', '-')),
        ('Turbines',        f"{n} × {d.get('wtg_model', '-')}"),
        ('Rotor diameter',  f"{d.get('rotor_m', '-')} m"),
        ('Closest WTG pair', f"{min_sp:.1f}D" if min_sp else '-'),
        ('Hub height',      f"{int(d.get('hub_m', 0) or 0)} m"),
        ('Total capacity',  f"{mw} MW"),
        ('Gross AEP',       f"{d.get('gross_aep_mwh', 0)/1000:.1f} GWh"),
        ('Wake loss',       f"{d.get('wake_loss_pct', 0):.1f}%"),
        ('Park AEP',        f"{d.get('park_yield_mwh', 0)/1000:.1f} GWh"),
        ('P50 AEP',         f"{p50/1000:.1f} GWh"),
        ('Capacity factor', cf),
    ]
    highlight = {'Gross AEP', 'P50 AEP'}

    _mini_label(slide, 'Results Summary',
                Inches(cx0), Inches(cy0), Inches(left_w))

    tbl = slide.shapes.add_table(
        len(metrics), 2,
        Inches(cx0), Inches(cy0 + lbl_h),
        Inches(left_w), Inches(tbl_h)
    ).table
    half = int(Inches(left_w) * 0.50)
    tbl.columns[0].width = half
    tbl.columns[1].width = int(Inches(left_w)) - half

    for ri, (label, value) in enumerate(metrics):
        hl     = label in highlight
        lbl_fg = ORANGE if hl else WHITE
        dat_fg = ORANGE if hl else BLACK
        _set_cell(tbl.cell(ri, 0), label, bg=DARK_NAVY, fg=lbl_fg,
                  bold=hl, size=8, align=PP_ALIGN.LEFT)
        _set_cell(tbl.cell(ri, 1), value, bg=WHITE, fg=dat_fg,
                  bold=hl, size=8)
        tbl.rows[ri].height = Inches(0.164)

    # ── Left panel: wake chart ────────────────────────────────────────────
    wtgs = d.get('wtgs', [])
    if wtgs and chart_h > 0.5:
        _mini_label(slide, 'Per-WTG Wake Loss [%]',
                    Inches(cx0), Inches(chart_lbl_y), Inches(left_w))

        chart_png = wake_chart_bytes(
            wtgs,
            avg_pct=d.get('wake_loss_pct', 0),
            width_in=left_w * 0.95,
            height_in=chart_h * 0.85)
        slide.shapes.add_picture(
            io.BytesIO(chart_png),
            Inches(cx0),
            Inches(chart_lbl_y + lbl_h),
            width=Inches(left_w),
            height=Inches(chart_h))


# ─────────────────────────────────────────────────────────────────────────────
# Summary comparison table slide
# ─────────────────────────────────────────────────────────────────────────────

def add_summary_slide(prs: Presentation, datasets: list[dict],
                      blank_layout, losses: dict):
    """Final slide: comparison table, one column per calculation."""
    cx0 = LM
    cy0 = HDR_H + 0.08
    cw  = RM - LM
    ch  = FOOT_Y - cy0

    slide = prs.slides.add_slide(blank_layout)
    _set_slide_title(slide, 'PARK Calculation Summary')

    def _gwh(key):
        return lambda d: f"{d[key]/1000:.1f}" if d.get(key) else '-'
    def _pct(key):
        return lambda d: f"{d[key]:.1f}%" if d.get(key) is not None else '-'
    def _plain(key):
        return lambda d: str(d[key]) if d.get(key) is not None else '-'
    def _loss(k):
        # Use per-PDF losses stored on each dataset (set by build()); fall back
        # to the global losses dict if not present.
        def fn(d):
            v = d.get('_losses', losses).get(k)
            return 'ignored' if v is None else f'{v:.1f}%'
        return fn

    rows = [
        ('Calculation date',                  _plain('calc_date'),      'normal'),
        ('Manufacturer / WTG name',           _plain('manufacturer'),   'bold'),
        ('Rotor diameter [m]',                _plain('rotor_m'),        'normal'),
        ('Closest WTG pair [D]',
            lambda d: f"{d['min_spacing_d']:.1f}D" if d.get('min_spacing_d') else '-',
            'normal'),
        ('Rated power [MW]',                  _plain('rated_mw'),       'normal'),
        ('Nameplate [MW]',                    _plain('total_mw'),       'orange'),
        ('Hub Height [m]',
            lambda d: str(int(d['hub_m'])) if d.get('hub_m') else '-', 'bold'),
        ('Number of WTGs',                    _plain('num_wtgs'),       'bold'),
        ("Mean free wind speed [m/s]",        _plain('mean_free_ws'),   'normal'),
        ("Mean reduced wind speed [m/s]",     _plain('mean_reduced_ws'),'normal'),
        ('Gross Yield [GWh]',                 _gwh('gross_aep_mwh'),    'bold'),
        ('Wake Losses [%]',                   _pct('wake_loss_pct'),    'bold'),
        ('Park Yield [GWh]  (wake only)',      _gwh('park_yield_mwh'),   'bold'),
        ('Loss Assumptions:',                 lambda _: '',             'section'),
        ('  Temp derating loss [%]',          _loss('Temp derating loss [%]'),          'indent_dark'),
        ('  Availability loss [%]',           _loss('Availability loss [%]'),           'indent'),
        ('  Electrical loss [%]',             _loss('Electrical loss [%]'),             'indent'),
        ('  Turbine performance loss [%]',    _loss('Turbine performance loss [%]'),    'indent'),
        ('  Degradation [%]',                 _loss('Degradation [%]'),                 'indent'),
        ('P50 AEP [GWh]  (incl. losses)',     _gwh('p50_aep_mwh'),     'orange'),
        ('AEP difference vs. baseline [%]',   None,                    'bold'),
        ('Energy per turbine [GWh/yr]',
            lambda d: f"{d['p50_aep_mwh']/1000/d['num_wtgs']:.1f}"
                if d.get('p50_aep_mwh') and d.get('num_wtgs') else '-', 'bold'),
        ('Capacity factor (after losses) [%]',
            lambda d: f"{d['p50_aep_mwh']/(d['total_mw']*8760)*100:.1f}%"
                if d.get('p50_aep_mwh') and d.get('total_mw') else '-', 'normal'),
    ]

    n_cols = 1 + len(datasets)
    n_rows = 1 + len(rows)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(cx0), Inches(cy0),
        Inches(cw), Inches(ch)
    ).table

    label_w = int(Inches(cw) * 0.36)
    data_w  = (int(Inches(cw)) - label_w) // len(datasets)
    tbl.columns[0].width = label_w
    for i in range(1, n_cols):
        tbl.columns[i].width = data_w

    # Header row — dark navy throughout
    _set_cell(tbl.cell(0, 0), '', bg=DARK_NAVY)
    for ci, d in enumerate(datasets):
        _set_cell(tbl.cell(0, ci + 1),
                  d.get('calc_name', f'Calc {ci+1}'),
                  bg=DARK_NAVY, fg=WHITE, bold=True, size=6.5)
    tbl.rows[0].height = Inches(0.164)

    base_p50 = datasets[0].get('p50_aep_mwh') if datasets else None

    for ri, (label, val_fn, style) in enumerate(rows):
        is_orange  = style == 'orange'
        is_indent  = style in ('indent', 'indent_dark')
        is_bold    = style in ('bold', 'orange', 'section')
        font_size  = 5.5 if is_indent else 6.5

        # Label column: always dark navy bg
        lbl_fg = ORANGE if is_orange else WHITE
        _set_cell(tbl.cell(ri + 1, 0), label.lstrip(),
                  bg=DARK_NAVY, fg=lbl_fg,
                  bold=is_bold and not is_indent,
                  italic=is_indent, size=font_size, align=PP_ALIGN.LEFT)

        # Data columns: always white bg; grey text for loss/indent rows
        dat_fg = ORANGE if is_orange else (GREY if is_indent else BLACK)

        # 'indent' rows (fixed losses) — same value for every calc → merge cells
        if style == 'indent' and len(datasets) > 1:
            const_val = val_fn(datasets[0]) if val_fn else '-'
            for ci in range(len(datasets)):
                _set_cell(tbl.cell(ri + 1, ci + 1),
                          const_val if ci == 0 else '',
                          bg=WHITE, fg=dat_fg,
                          bold=False, italic=True, size=font_size)
            _hmerge_row(tbl, ri + 1, 1, len(datasets))
        else:
            for ci, d in enumerate(datasets):
                if label == 'AEP difference vs. baseline [%]':
                    if ci == 0 or base_p50 is None:
                        val = '-'
                    elif d.get('p50_aep_mwh'):
                        diff = (d['p50_aep_mwh'] - base_p50) / base_p50 * 100
                        val  = f'{diff:+.1f}%'
                    else:
                        val = '-'
                else:
                    val = val_fn(d) if val_fn else '-'

                _set_cell(tbl.cell(ri + 1, ci + 1), val,
                          bg=WHITE, fg=dat_fg,
                          bold=is_bold and not is_indent,
                          italic=is_indent, size=font_size)

        tbl.rows[ri + 1].height = Inches(0.139 if is_indent else 0.164)


# ─────────────────────────────────────────────────────────────────────────────
# Public build() — called by both CLI and Streamlit app
# ─────────────────────────────────────────────────────────────────────────────

def build(pdf_paths: list[str],
          template_path: str | None = None,
          losses: dict | None = None,
          cover_title: str = 'XXWF Prelim Yield Estimates',
          cover_subtitle: str = 'Version Y',
          cover_subsubtitle: str = '',
          losses_per_pdf: list[dict] | None = None,
          shapes: list | None = None,
          shapes_per_calc: list[list] | None = None,
          noise_overlays: list | None = None) -> bytes:
    """
    Extract data from each PDF, build the presentation, return .pptx bytes.

    Parameters
    ----------
    pdf_paths      : ordered list of WindPRO PARK PDF paths
    template_path  : .pptx template (slide master / layouts preserved)
    losses         : default loss assumptions dict; defaults to DEFAULT_LOSSES
    cover_title    : title text for the cover slide
    cover_subtitle    : heading 2 on cover (orange)
    cover_subsubtitle : heading 3 on cover (orange, smaller)
    losses_per_pdf : optional per-PDF loss dicts (same length as pdf_paths);
                     when supplied, each PDF uses its own loss assumptions
    """
    if losses is None:
        losses = DEFAULT_LOSSES

    datasets, pdf_strs = [], []
    for i, p in enumerate(pdf_paths):
        data = extract(str(p))
        per = losses_per_pdf[i] if losses_per_pdf else losses
        apply_losses(data, per)
        data['_losses'] = per   # carried through to summary slide
        datasets.append(data)
        pdf_strs.append(str(p))

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        # Remove template example slides (keep slide master / layouts).
        # Must drop from _rels too, else python-pptx writes orphaned parts
        # to the ZIP causing duplicate slide XML entries.
        sldIdLst = prs.slides._sldIdLst
        for slide in list(prs.slides):
            slide_part = slide.part
            for rId, rel in prs.part._rels.items():
                if getattr(rel, '_target', None) is slide_part:
                    prs.part._rels._rels.pop(rId, None)
                    break
        for el in list(sldIdLst):
            sldIdLst.remove(el)
    else:
        prs = Presentation()
        prs.slide_width  = Inches(10)
        prs.slide_height = Inches(5.63)

    blank_layout = _layout_by_name(prs, 'Content_Blank')

    # Cover slide
    add_cover_slide(prs, cover_title, cover_subtitle, cover_subsubtitle)

    # Per-calculation slides
    for i, (data, pdf_str) in enumerate(zip(datasets, pdf_strs)):
        per_shapes  = shapes_per_calc[i]  if shapes_per_calc  else shapes
        per_noise   = noise_overlays[i]   if noise_overlays   else None
        add_calc_slide(prs, data, pdf_str, blank_layout, losses,
                       shapes=per_shapes, noise_overlay=per_noise)

    # Summary table slide (needs at least one dataset — its column widths
    # divide by the calculation count)
    if datasets:
        add_summary_slide(prs, datasets, blank_layout, losses)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# CLI entry point
# ─────────────────────────────────────────────────────────────────────────────

def main(folder=None):
    # folder: where to look for the default template / PARK_*.pdf and where to
    # write the output deck. The tool-repo shims pass their own directory; a
    # direct invocation falls back to the current working directory.
    args = sys.argv[1:]
    folder = Path(folder) if folder is not None else Path.cwd()

    template_path = None
    if '--template' in args:
        idx = args.index('--template')
        template_path = args[idx + 1]
        args = args[:idx] + args[idx + 2:]

    if template_path is None:
        default_tpl = folder / 'template01.pptx'
        if default_tpl.exists():
            template_path = str(default_tpl)
            print(f'Using default template: {default_tpl.name}')

    pdf_paths = [str(Path(p)) for p in args] if args \
        else [str(p) for p in sorted(folder.glob('PARK_*.pdf'))]

    if not pdf_paths:
        print('No PARK_*.pdf files found. Pass PDF paths as arguments.')
        sys.exit(1)

    print(f'Processing {len(pdf_paths)} PDF(s)…')
    for p in pdf_paths:
        d = extract(p)
        apply_losses(d, DEFAULT_LOSSES)
        n_coords = len(d.get('wtg_coords', {}))
        print(f'  {Path(p).name}')
        print(f'    {len(d["wtgs"])} WTGs | '
              f'Gross {d.get("gross_aep_mwh",0)/1000:.1f} GWh | '
              f'Wake {d.get("wake_loss_pct",0):.1f}% | '
              f'P50 {d.get("p50_aep_mwh",0)/1000:.1f} GWh | '
              f'{n_coords} coords')

    from datetime import date
    today = date.today()
    subsubtitle = f"{_ordinal(today.day)} {today.strftime('%B %Y')}"
    pptx_bytes = build(pdf_paths, template_path, cover_subsubtitle=subsubtitle)

    out = folder / 'park_summary.pptx'
    out.write_bytes(pptx_bytes)
    print(f'\nSaved → {out}')


if __name__ == '__main__':
    main()
